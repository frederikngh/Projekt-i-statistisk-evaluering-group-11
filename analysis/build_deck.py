"""Builds statistical_tests.pptx - short midway check-in on the tests used."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

RED = RGBColor(0x99, 0x00, 0x00)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x5A, 0x5A, 0x5A)
LIGHT = RGBColor(0xF2, 0xEF, 0xEC)
HEAD = RGBColor(0x77, 0x00, 0x00)
GREEN = RGBColor(0x2E, 0x6E, 0x33)
AMBER = RGBColor(0xA8, 0x5A, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SHOT = "data/screenshots/Fall2018_Q10.png"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
blank = prs.slide_layouts[6]


def textbox(slide, left, top, width, height, lines, anchor=MSO_ANCHOR.TOP, mono=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, ln in enumerate(lines):
        text, size, bold, color = ln[0], ln[1], ln[2], ln[3]
        align = ln[4] if len(ln) > 4 else PP_ALIGN.LEFT
        space = ln[5] if len(ln) > 5 else 6
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space)
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = "Consolas" if mono else "Calibri"
    return tb


def title_bar(slide, kicker, title):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(1.25))
    bar.fill.solid(); bar.fill.fore_color.rgb = RED
    bar.line.fill.background(); bar.shadow.inherit = False
    textbox(slide, Inches(0.55), Inches(0.12), Inches(12.2), Inches(1.05),
            [(kicker, 13, True, RGBColor(0xF0, 0xC8, 0xC8)),
             (title, 26, True, WHITE)])


def panel(slide, left, top, width, height, edge=RED):
    p = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    p.fill.solid(); p.fill.fore_color.rgb = LIGHT
    p.line.color.rgb = edge; p.shadow.inherit = False
    return p


def set_cell(cell, text, size=14, bold=False, color=DARK, bg=WHITE, align=PP_ALIGN.LEFT):
    cell.fill.solid(); cell.fill.fore_color.rgb = bg
    cell.margin_left = Inches(0.12); cell.margin_right = Inches(0.08)
    cell.margin_top = Inches(0.03); cell.margin_bottom = Inches(0.03)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf = cell.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = "Calibri"


def table(slide, left, top, width, col_w, rows, row_h=0.5, head_h=0.55):
    nrows, ncols = len(rows), len(rows[0])
    gr = slide.shapes.add_table(nrows, ncols, left, top, width,
                                Inches(head_h + row_h * (nrows - 1)))
    t = gr.table
    t.first_row = False; t.horz_banding = False
    for j, w in enumerate(col_w):
        t.columns[j].width = Inches(w)
    t.rows[0].height = Inches(head_h)
    for i in range(1, nrows):
        t.rows[i].height = Inches(row_h)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            if i == 0:
                set_cell(t.cell(i, j), val, 14, True, WHITE, HEAD,
                         PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)
            else:
                bg = LIGHT if i % 2 == 0 else WHITE
                set_cell(t.cell(i, j), val, 14, False, DARK, bg,
                         PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)
    return gr


# ---------- 1. TITLE ----------
s = prs.slides.add_slide(blank)
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.3), SW, Inches(2.8))
band.fill.solid(); band.fill.fore_color.rgb = RED
band.line.fill.background(); band.shadow.inherit = False
textbox(s, Inches(0.7), Inches(2.45), Inches(12), Inches(2.5),
        [("02445 - Statistical Evaluation of AI   |   midway check-in", 16, True, RGBColor(0xF0, 0xC8, 0xC8)),
         ("Does Gemma do worse on graphs than on text?", 32, True, WHITE),
         ("The statistical tests we use - are they appropriate?", 19, False, RGBColor(0xF5, 0xE6, 0xE6))])
textbox(s, Inches(0.7), Inches(5.45), Inches(12), Inches(1.4),
        [("Gemma 4 answers 405 DTU ML exam multiple-choice questions (options A-D, plus E = don't know).",
          15, False, GRAY),
         ("Outcome is binary: correct / wrong, scored against the official answer key.", 15, False, GRAY)])

# ---------- 2. EXAMPLE: WHAT WE GIVE GEMMA ----------
s = prs.slides.add_slide(blank)
title_bar(s, "WHAT WE GIVE GEMMA", "The same question, two formats")
textbox(s, Inches(0.55), Inches(1.4), Inches(12.2), Inches(1.4),
        [("Consider the split in Table 3. We build a classification tree with only this split and evaluate it "
          "on the same data it was trained on. What is the accuracy?", 16, True, DARK),
         ("A. 0.64     B. 0.29     C. 0.35     D. 0.43     E. Don't know", 15, False, GRAY)])
# left panel: as text
panel(s, Inches(0.55), Inches(2.95), Inches(6.0), Inches(3.0))
textbox(s, Inches(0.8), Inches(3.1), Inches(5.5), Inches(0.4),
        [("GIVEN AS TEXT", 14, True, RED)])
textbox(s, Inches(0.8), Inches(3.6), Inches(5.5), Inches(2.2),
        [("x9-interval    y=1   y=2   y=3", 13, False, DARK, PP_ALIGN.LEFT, 4),
         ("x9 <= 0.13     108   112    56", 13, False, DARK, PP_ALIGN.LEFT, 4),
         ("0.13 < x9       58    75   116", 13, False, DARK, PP_ALIGN.LEFT, 4)], mono=True)
# right panel: as image
panel(s, Inches(6.85), Inches(2.95), Inches(5.9), Inches(3.0))
textbox(s, Inches(7.1), Inches(3.1), Inches(5.4), Inches(0.4),
        [("GIVEN AS IMAGE (cropped from the exam PDF)", 14, True, RED)])
pic = s.shapes.add_picture(SHOT, Inches(7.45), Inches(3.65), width=Inches(4.7))
pic.line.color.rgb = GRAY; pic.line.width = Pt(0.75)
textbox(s, Inches(0.55), Inches(6.15), Inches(12.2), Inches(1.1),
        [("Only the format of the data changes - the stem and options are identical in both arms.",
          15, True, DARK),
         ("Gemma replies with one letter; we score correct / wrong (here the key is D).", 14, False, GRAY)])

# ---------- 3. ACCURACY NUMBERS ----------
s = prs.slides.add_slide(blank)
title_bar(s, "THE DATA BEHIND THE TESTS", "Accuracy numbers we feed into the tests")
textbox(s, Inches(0.55), Inches(1.4), Inches(7.0), Inches(0.4),
        [("Accuracy by format", 16, True, DARK)])
table(s, Inches(0.55), Inches(1.85), Inches(7.0), [3.5, 1.9, 1.6],
      [["Format", "Correct / n", "Accuracy"],
       ["pure text", "98 / 134", "73.1%"],
       ["figure/table as text", "80 / 127", "63.0%"],
       ["figure/table as image", "54 / 127", "42.5%"],
       ["geometric figure (image only)", "52 / 144", "36.1%"]],
      row_h=0.58)
textbox(s, Inches(8.0), Inches(1.4), Inches(5.0), Inches(0.4),
        [("Paired set, image vs text (McNemar)", 16, True, DARK)])
table(s, Inches(8.0), Inches(1.85), Inches(4.8), [1.7, 1.55, 1.55],
      [["", "text right", "text wrong"],
       ["image right", "45", "9"],
       ["image wrong", "35", "38"]],
      row_h=0.7, head_h=0.55)
textbox(s, Inches(8.0), Inches(4.6), Inches(4.8), Inches(1.0),
        [("127 paired questions; only the 44", 13, False, GRAY),
         ("disagreements (35 vs 9) drive the test.", 13, False, GRAY)])
line = panel(s, Inches(0.55), Inches(6.0), Inches(12.25), Inches(0.95), edge=GRAY)
textbox(s, Inches(0.8), Inches(6.12), Inches(11.8), Inches(0.8),
        [("Contamination check (exams Gemma may have seen vs. not):  pre-cutoff 263 / 498 = 52.8%   "
          "vs.   post-cutoff 46 / 68 = 67.6%", 15, True, DARK, PP_ALIGN.LEFT, 0)],
        anchor=MSO_ANCHOR.MIDDLE)

# ---------- 4. THE TESTS ----------
s = prs.slides.add_slide(blank)
title_bar(s, "THE TESTS WE USE", "Each test and why it fits a binary outcome")
table(s, Inches(0.45), Inches(1.55), Inches(12.45), [3.5, 4.6, 4.35],
      [["Test", "What it compares", "Why it is appropriate"],
       ["Binomial test (exact)", "each format's accuracy vs. 25% chance", "exact test for one proportion; no large-n assumption"],
       ["McNemar, exact (paired) - PRIMARY", "same question as image vs. as text", "the paired binary test; pairing removes question difficulty"],
       ["Two-proportion z (unpaired)", "text questions vs. pure-graph questions", "standard 2-proportion test; n > 130 per group"],
       ["One-sided z", "pre-cutoff vs. post-cutoff exams", "2-proportion test; one-sided (contamination has a direction)"],
       ["95% confidence interval", "every accuracy and difference", "course-standard normal interval; paired diff uses paired formula"]],
      row_h=0.92, head_h=0.5)
textbox(s, Inches(0.45), Inches(6.85), Inches(12.4), Inches(0.5),
        [("Binary outcome -> the proportion / McNemar family (not t-tests or ANOVA, which need a numeric outcome).",
          14, True, DARK)])

# ---------- 5. FOR YOUR FEEDBACK ----------
s = prs.slides.add_slide(blank)
title_bar(s, "WHAT WE WOULD LIKE YOUR INPUT ON", "Design choices to confirm")
textbox(s, Inches(0.7), Inches(1.7), Inches(12), Inches(5.2),
        [("1.  Binary outcome (correct/wrong) -> proportion / McNemar family rather than t-test / ANOVA. OK?",
          18, False, DARK, PP_ALIGN.LEFT, 14),
         ("2.  McNemar is our single confirmatory (primary) test -> we argue no multiple-testing correction is needed.",
          18, False, DARK, PP_ALIGN.LEFT, 14),
         ("3.  The unpaired text-vs-graph z has difficulty confounded with format -> we report it as secondary.",
          18, False, DARK, PP_ALIGN.LEFT, 14),
         ("4.  Contamination test is one-sided (pre-cutoff > post-cutoff). Is one-sided the right call?",
          18, False, DARK, PP_ALIGN.LEFT, 14),
         ("5.  Confidence intervals use the 02402 normal approximation (we dropped Wilson / bootstrap as off-syllabus).",
          18, False, DARK, PP_ALIGN.LEFT, 14),
         ("", 8, False, DARK, PP_ALIGN.LEFT, 6),
         ("Are these five choices appropriate, or would you change any before we scale up the analysis?",
          17, True, RED, PP_ALIGN.LEFT, 0)])

prs.save("statistical_tests.pptx")
print("saved with", len(prs.slides._sldIdLst), "slides")
